from __future__ import annotations

import asyncio
import hashlib
import json
import mimetypes
import re
from pathlib import Path
from typing import Any

from fastapi import UploadFile

from local_api.config import get_settings
from local_api.db import get_store, new_id, utc_now
from local_api.services.coursebuilder import get_coursebuilder_service
from local_api.services.knowledge_graph import get_knowledge_graph_service
from local_api.services.settings import get_settings_service
from local_api.services.vector_service import get_vector_service


class IngestionService:
    def __init__(self) -> None:
        self._tasks: set[asyncio.Task[Any]] = set()

    async def ingest_upload(self, conversation_id: str, upload: UploadFile) -> dict:
        raw = await upload.read()
        filename = upload.filename or "upload.txt"
        mime_type = upload.content_type or mimetypes.guess_type(filename)[0] or "application/octet-stream"
        if not raw:
            raise ValueError("empty upload")
        file_hash = hashlib.sha256(raw).hexdigest()[:16]
        safe_name = re.sub(r"[^A-Za-z0-9._-]+", "_", filename).strip("_") or "upload.bin"
        stored_path = get_settings().data_dir / "objects" / "uploads" / f"{file_hash}_{safe_name}"
        stored_path.parent.mkdir(parents=True, exist_ok=True)
        stored_path.write_bytes(raw)
        record = get_store().create_uploaded_file(
            conversation_id,
            safe_name,
            stored_path,
            mime_type=mime_type,
            size_bytes=len(raw),
        )
        get_coursebuilder_service().invalidate_plan(conversation_id, "source file added")
        return record

    async def retry_upload(self, conversation_id: str, file_id: str) -> dict:
        record = get_store().get_file_for_conversation(conversation_id, file_id)
        if record is None:
            raise FileNotFoundError("file not found")
        if record["status"] != "failed":
            raise ValueError("only failed files can be retried")
        get_store().clear_file_content(conversation_id, file_id)
        get_coursebuilder_service().invalidate_plan(conversation_id, "source file retry requested")
        get_store().update_file(
            file_id,
            status="uploaded",
            parser_used=None,
            error=None,
            chunk_count=0,
        )
        return get_store().get_file(file_id) or record

    async def process_upload(self, file_id: str) -> dict:
        record = get_store().get_file(file_id)
        if record is None:
            return {"ok": False, "error": "file not found"}
        return await self._parse_and_index(record)

    async def process_upload_batch(self, file_ids: list[str]) -> list[dict]:
        return list(await asyncio.gather(*(self.process_upload(file_id) for file_id in file_ids)))

    def resume_incomplete_uploads(self) -> None:
        rows = get_store().query(
            "SELECT id FROM uploaded_files WHERE status NOT IN ('ready', 'failed') ORDER BY created_at ASC"
        )
        for row in rows:
            task = asyncio.create_task(self.process_upload(row["id"]))
            self._tasks.add(task)
            task.add_done_callback(self._tasks.discard)

    async def _parse_and_index(self, file_record: dict) -> dict:
        file_id = file_record["id"]
        conversation_id = file_record["conversation_id"]
        source_path = Path(file_record["stored_path"])
        get_coursebuilder_service().invalidate_plan(conversation_id, "source file processing restarted")
        get_store().clear_file_content(conversation_id, file_id)
        get_store().update_file(file_id, status="parsing", error=None, chunk_count=0)
        try:
            text, parser_used = await self._parse_file(source_path, file_record["filename"])
            if not text.strip():
                raise ValueError("no extractable text")
            parsed_key = get_settings().data_dir / "objects" / "parsed" / f"{file_id}.md"
            cleaned_key = get_settings().data_dir / "objects" / "cleaned" / f"{file_id}.txt"
            parsed_key.parent.mkdir(parents=True, exist_ok=True)
            cleaned_key.parent.mkdir(parents=True, exist_ok=True)
            parsed_key.write_text(text, encoding="utf-8")
            cleaned = _clean_text(text)
            cleaned_key.write_text(cleaned, encoding="utf-8")
            get_store().update_file(file_id, status="chunking", parser_used=parser_used)
            document_id = self._insert_document(conversation_id, file_id, file_record["filename"])
            sections = _split_sections(cleaned, file_record["filename"])
            chunks = []
            chunk_index = 0
            for section_index, section in enumerate(sections):
                section_id = self._insert_section(
                    conversation_id,
                    document_id,
                    file_id,
                    section["heading_path"],
                    section["text"],
                    section_index,
                )
                section_chunk_index = 0
                for chunk_text in _chunk_text(section["text"], get_settings().chunk_max_chars, get_settings().chunk_overlap_chars):
                    concepts = _extract_concepts(chunk_text)
                    questions = _generated_questions(concepts, section["heading_path"], get_settings().generated_questions_per_chunk)
                    equations = _extract_equations(chunk_text)
                    tables = _extract_tables(chunk_text)
                    timeline_events = _extract_timeline_events(chunk_text)
                    chunks.append(
                        {
                            "id": _stable_chunk_id(file_id, section_id, section_chunk_index, chunk_text),
                            "conversation_id": conversation_id,
                            "source_file_id": file_id,
                            "document_id": document_id,
                            "section_id": section_id,
                            "chunk_index": chunk_index,
                            "text": chunk_text,
                            "source_filename": file_record["filename"],
                            "metadata": {
                                "chunker": "structured-section-v1",
                                "document_id": document_id,
                                "section_id": section_id,
                                "parent_section_id": "",
                                "source_file_id": file_id,
                                "source": file_record["filename"],
                                "heading_path": " > ".join(section["heading_path"]),
                                "heading_path_list": section["heading_path"],
                                "section_title": section["heading_path"][-1] if section["heading_path"] else file_record["filename"],
                                "section_index": section_index,
                                "chunk_index": chunk_index,
                                "section_chunk_index": section_chunk_index,
                                "key_concepts": concepts,
                                "generated_questions": questions,
                                "context_type": "focused_chunk",
                                "section_summary": _summary(section["text"]),
                                "equation_count": len(equations),
                                "table_count": len(tables),
                                "timeline_event_count": len(timeline_events),
                                "equations": equations,
                                "tables": tables,
                                "timeline_events": timeline_events,
                            },
                        }
                    )
                    chunk_index += 1
                    section_chunk_index += 1
            _link_neighbor_chunks(chunks)
            get_store().update_file(file_id, status="extracting_concepts", chunk_count=len(chunks))
            get_store().replace_chunks_for_file(file_id, chunks)
            self._store_content_features(conversation_id, file_id, chunks)
            self._upsert_concepts(conversation_id, chunks)
            get_store().update_file(file_id, status="planning_course", chunk_count=len(chunks))
            plan = await get_coursebuilder_service().prepare_plan_async(conversation_id)
            if plan.get("status") not in {"draft", "validated"}:
                raise ValueError("course plan could not be prepared before embedding")
            get_store().update_file(file_id, status="embedding", chunk_count=len(chunks))
            embeddings = await get_vector_service().embed_chunks(chunks)
            retrieval_settings = get_settings_service().effective_retrieval_settings()
            get_store().update_chunk_embeddings(
                embeddings,
                model=retrieval_settings.embedding_model,
                dim=retrieval_settings.embedding_dim,
            )
            get_store().update_file(file_id, status="building_course", chunk_count=len(chunks))
            get_knowledge_graph_service().rebuild_graph(conversation_id)
            get_store().update_file(file_id, status="ready", parser_used=parser_used, error=None, chunk_count=len(chunks))
            try:
                ready_files = get_store().list_files(conversation_id)
                if ready_files and all(file["status"] == "ready" for file in ready_files):
                    await get_coursebuilder_service().rebuild_async(conversation_id)
            except Exception:
                pass
        except Exception as exc:  # noqa: BLE001
            get_store().update_file(file_id, status="failed", error=str(exc))
        return get_store().get_file(file_id) or file_record

    async def _parse_file(self, path: Path, filename: str) -> tuple[str, str]:
        parser_settings = get_settings_service().get_parser_settings()
        if not parser_settings.use_local_parsers_only:
            api_key = get_settings_service().parser_api_key()
            if not api_key:
                raise ValueError("LlamaCloud API key is required when LlamaParse mode is enabled")
            return await _parse_with_llama_cloud(path, filename, api_key)
        return _parse_local_file(path, filename)

    def _insert_document(self, conversation_id: str, file_id: str, title: str) -> str:
        document_id = new_id("doc")
        get_store().execute(
            """
            INSERT INTO course_documents (id, conversation_id, source_file_id, title, metadata_json, created_at)
            VALUES (?, ?, ?, ?, '{}', ?)
            """,
            (document_id, conversation_id, file_id, title, utc_now()),
        )
        return document_id

    def _insert_section(
        self,
        conversation_id: str,
        document_id: str,
        file_id: str,
        heading_path: list[str],
        text: str,
        order_index: int,
    ) -> str:
        section_id = new_id("section")
        summary = text[:320].rsplit(" ", 1)[0]
        import json

        get_store().execute(
            """
            INSERT INTO course_sections
              (id, conversation_id, document_id, source_file_id, heading_path_json, summary, order_index, text, created_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                section_id,
                conversation_id,
                document_id,
                file_id,
                json.dumps(heading_path),
                summary,
                order_index,
                text,
                utc_now(),
            ),
        )
        return section_id

    def _upsert_concepts(self, conversation_id: str, chunks: list[dict]) -> None:
        seen: set[str] = set()
        for chunk in chunks:
            for concept in chunk.get("metadata", {}).get("key_concepts", []):
                key = concept.casefold()
                if key in seen:
                    continue
                seen.add(key)
                get_store().execute(
                    """
                    INSERT INTO concept_inventory (id, conversation_id, name, aliases_json, metadata_json)
                    VALUES (?, ?, ?, '[]', '{}')
                    """,
                    (new_id("concept"), conversation_id, concept),
                )

    def _store_content_features(self, conversation_id: str, file_id: str, chunks: list[dict]) -> None:
        for chunk in chunks:
            metadata = chunk.get("metadata", {})
            for expression in metadata.get("equations", []):
                kind = "matrix" if "\\begin{" in expression and "matrix}" in expression else (
                    "chemical" if _looks_like_chemical_equation(expression) else "equation"
                )
                get_store().execute(
                    """
                    INSERT INTO formulas (id, conversation_id, source_file_id, chunk_id, label, expression, metadata_json)
                    VALUES (?, ?, ?, ?, ?, ?, ?)
                    """,
                    (
                        new_id("formula"), conversation_id, file_id, chunk["id"], kind, expression,
                        json.dumps({"kind": kind}, ensure_ascii=False),
                    ),
                )
            for table_markdown in metadata.get("tables", []):
                rows = [
                    [cell.strip() for cell in line.strip().strip("|").split("|")]
                    for line in table_markdown.splitlines()
                    if line.strip() and not re.match(r"^\s*\|?\s*:?-{3,}", line)
                ]
                get_store().execute(
                    """
                    INSERT INTO tables (id, conversation_id, source_file_id, chunk_id, caption, content_json, metadata_json)
                    VALUES (?, ?, ?, ?, ?, ?, '{}')
                    """,
                    (
                        new_id("table"), conversation_id, file_id, chunk["id"], "Source table",
                        json.dumps({"headers": rows[0] if rows else [], "rows": rows[1:]}, ensure_ascii=False),
                    ),
                )
            for event in metadata.get("timeline_events", []):
                date_match = re.search(r"\b(?:\d{3,4}(?:\s*(?:BCE|BC|CE|AD))?|Q[1-4]\s+\d{4})\b", event, re.IGNORECASE)
                get_store().execute(
                    """
                    INSERT INTO timeline_events
                      (id, conversation_id, source_file_id, chunk_id, label, event_date, metadata_json)
                    VALUES (?, ?, ?, ?, ?, ?, '{}')
                    """,
                    (new_id("event"), conversation_id, file_id, chunk["id"], event, date_match.group(0) if date_match else None),
                )


def _parse_local_file(path: Path, filename: str) -> tuple[str, str]:
    suffix = Path(filename).suffix.lower()
    if suffix in {".txt", ".md", ".markdown", ".csv", ".tsv", ".json"}:
        return path.read_text(encoding="utf-8", errors="ignore"), "local:text"
    if suffix in {".html", ".htm"}:
        return _parse_html(path), "local:html"
    if suffix == ".pdf":
        return _parse_pdf(path), "local:pdf"
    if suffix == ".docx":
        return _parse_docx(path), "local:docx"
    if suffix == ".pptx":
        return _parse_pptx(path), "local:pptx"
    return path.read_text(encoding="utf-8", errors="ignore"), "local:raw_text"


async def _parse_with_llama_cloud(path: Path, filename: str, api_key: str) -> tuple[str, str]:
    try:
        from llama_cloud import AsyncClient
    except ImportError as exc:
        raise RuntimeError("llama-cloud>=1.0 is required for LlamaParse mode") from exc

    settings = get_settings()
    async with AsyncClient(api_key=api_key, base_url=settings.llama_cloud_base_url) as client:
        result = await client.parsing.parse(
            tier="cost_effective",
            version="latest",
            upload_file=(filename, path.read_bytes()),
            expand=["markdown"],
            polling_interval=settings.llama_cloud_poll_interval_s,
            timeout=settings.llama_cloud_timeout_s,
        )
    markdown = _extract_llama_markdown(result)
    job_id = _extract_llama_job_id(result)
    return markdown, f"llamacloud:{job_id or 'latest'}"


def _extract_llama_markdown(result: object) -> str:
    for attr in ("markdown_full", "markdown", "text_full", "text"):
        value = getattr(result, attr, None)
        if isinstance(value, str) and value.strip():
            return value
        pages = getattr(value, "pages", None)
        if pages:
            joined = "\n\n".join(
                page_markdown
                for page_markdown in (getattr(page, "markdown", None) for page in pages)
                if page_markdown
            )
            if joined.strip():
                return joined
    raise ValueError(f"LlamaCloud parse result had no markdown payload (got {result!r})")


def _extract_llama_job_id(result: object) -> str:
    job = getattr(result, "job", None)
    for holder in (job, result):
        for attr in ("id", "job_id"):
            value = getattr(holder, attr, None)
            if value:
                return str(value)
    return ""


def _parse_html(path: Path) -> str:
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
        return soup.get_text("\n")
    except Exception:  # noqa: BLE001
        return path.read_text(encoding="utf-8", errors="ignore")


def _parse_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        return "\n\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception as exc:  # noqa: BLE001
        raise ValueError(f"PDF local parser failed: {exc}") from exc


def _parse_docx(path: Path) -> str:
    try:
        import docx

        document = docx.Document(str(path))
        return "\n".join(paragraph.text for paragraph in document.paragraphs)
    except Exception as exc:  # noqa: BLE001
        raise ValueError(f"DOCX local parser failed: {exc}") from exc


def _parse_pptx(path: Path) -> str:
    try:
        from pptx import Presentation

        deck = Presentation(str(path))
        texts: list[str] = []
        for slide_index, slide in enumerate(deck.slides, start=1):
            texts.append(f"# Slide {slide_index}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n\n".join(texts)
    except Exception as exc:  # noqa: BLE001
        raise ValueError(f"PPTX local parser failed: {exc}") from exc


_HTML_COMMENT_RE = re.compile(r"<!--.*?-->", re.DOTALL)
_MARKDOWN_IMAGE_RE = re.compile(r"!\[[^\]]*]\([^)]*\)")
_HTML_IMAGE_RE = re.compile(r"<img\b[^>]*>", re.IGNORECASE)
_PAGE_COUNTER_RE = re.compile(r"\b\d{1,3}\s*/\s*\d{1,3}\b")
_DOT_LEADER_RE = re.compile(r"\.{5,}\s*\d{1,4}\s*$")
_DANGLING_IMAGE_FILE_RE = re.compile(r"\bpage_\d+_[\w_]+\.(?:png|jpg|jpeg|webp)\b", re.IGNORECASE)
_MULTISPACE_RE = re.compile(r"[ \t]{2,}")
_REPEATED_DOTS_RE = re.compile(r"(?:\s*\.\s*){5,}")
_MONTH_DATE_RE = re.compile(
    r"\b(?:january|february|march|april|may|june|july|august|september|october|november|december|"
    r"janvier|fevrier|mars|avril|mai|juin|juillet|aout|septembre|octobre|novembre|decembre)"
    r"\s+\d{1,2},?\s+\d{4}\b",
    re.IGNORECASE,
)
_EMPTY_PARENS_RE = re.compile(r"\(\s*[,;|]?\s*\)")
_ORG_FOOTER_RE = re.compile(
    r"\b(?:university|universite|faculty|faculte|school|college|institute|department|dept\.?|"
    r"professor|enseignant|copyright|all rights reserved)\b",
    re.IGNORECASE,
)
_NOISE_PHRASES = (
    "navigation controls",
    "navigation icons",
    "footer icons",
    "footer navigation",
    "page number icon",
    "presentation controls",
    "slide navigation",
    "seal or stamp",
    "stamp/seal",
    "logo",
)


def _clean_text(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = _HTML_COMMENT_RE.sub(" ", text)
    text = _MARKDOWN_IMAGE_RE.sub(" ", text)
    text = _HTML_IMAGE_RE.sub(" ", text)

    out: list[str] = []
    blank_pending = False
    previous_normalized = ""
    for raw_line in text.splitlines():
        line = _clean_line(raw_line)
        if not line:
            blank_pending = bool(out)
            continue
        normalized = " ".join(line.lower().split())
        if _should_drop_line(line, normalized):
            blank_pending = bool(out)
            continue
        if normalized == previous_normalized:
            continue
        if blank_pending and out and out[-1] != "":
            out.append("")
        out.append(line)
        previous_normalized = normalized
        blank_pending = False
    return "\n".join(out).strip()


def _clean_line(raw_line: str) -> str:
    line = raw_line.strip()
    if not line:
        return ""
    if line.startswith("|") and line.endswith("|"):
        return line
    line = _EMPTY_PARENS_RE.sub(" ", line)
    line = _MONTH_DATE_RE.sub(" ", line)
    line = _PAGE_COUNTER_RE.sub(" ", line)
    line = _DANGLING_IMAGE_FILE_RE.sub(" ", line)
    line = _REPEATED_DOTS_RE.sub(" ", line)
    line = _MULTISPACE_RE.sub(" ", line)
    return line.strip(" |\t")


def _should_drop_line(line: str, normalized: str) -> bool:
    if not normalized:
        return True
    stripped = line.strip()
    if stripped.startswith("|") and stripped.endswith("|"):
        return False
    if any(phrase in normalized for phrase in _NOISE_PHRASES):
        return True
    if _DOT_LEADER_RE.search(line):
        return True
    if _PAGE_COUNTER_RE.fullmatch(line):
        return True
    if _MONTH_DATE_RE.fullmatch(line):
        return True
    if _DANGLING_IMAGE_FILE_RE.search(line):
        return True
    if _is_likely_footer_line(line, normalized):
        return True
    if _is_mostly_punctuation(line):
        return True
    return False


def _is_likely_footer_line(line: str, normalized: str) -> bool:
    if len(line) > 180:
        return False
    has_footer_cue = bool(_ORG_FOOTER_RE.search(line))
    has_counter_or_date = bool(_PAGE_COUNTER_RE.search(line) or _MONTH_DATE_RE.search(line))
    separator_count = line.count("|") + line.count(" - ") + line.count(" / ")
    if has_footer_cue and (has_counter_or_date or separator_count >= 2):
        return True
    return has_footer_cue and len(normalized.split()) <= 10


def _is_mostly_punctuation(line: str) -> bool:
    if len(line) < 8:
        return False
    useful = sum(1 for ch in line if ch.isalnum())
    return useful / len(line) < 0.25


def _split_sections(text: str, fallback_title: str) -> list[dict]:
    sections: list[dict] = []
    current_heading = [fallback_title]
    heading_stack: list[str] = []
    current: list[str] = []
    heading_re = re.compile(r"^(#{1,6})\s+(.+)$")
    for line in text.splitlines():
        match = heading_re.match(line.strip())
        if match and current:
            sections.append({"heading_path": current_heading, "text": "\n".join(current).strip()})
            current = []
            level = len(match.group(1))
            heading_stack = heading_stack[: level - 1] + [match.group(2).strip()]
            current_heading = heading_stack or [fallback_title]
        elif match:
            level = len(match.group(1))
            heading_stack = heading_stack[: level - 1] + [match.group(2).strip()]
            current_heading = heading_stack or [fallback_title]
        else:
            current.append(line)
    if current:
        sections.append({"heading_path": current_heading, "text": "\n".join(current).strip()})
    if not sections:
        sections.append({"heading_path": [fallback_title], "text": text})
    return [section for section in sections if section["text"].strip()]


def _chunk_text(text: str, max_chars: int, overlap_chars: int) -> list[str]:
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    chunks: list[str] = []
    current = ""
    for paragraph in paragraphs:
        if len(current) + len(paragraph) + 2 <= max_chars:
            current = f"{current}\n\n{paragraph}".strip()
            continue
        if current:
            chunks.append(current)
        prefix = current[-overlap_chars:] if current and overlap_chars > 0 else ""
        current = f"{prefix}\n\n{paragraph}".strip()
    if current:
        chunks.append(current[:max_chars])
    return chunks or [text[:max_chars]]


def _stable_chunk_id(source_file_id: str, section_id: str, section_chunk_index: int, text: str) -> str:
    fingerprint = hashlib.sha1(text.encode("utf-8")).hexdigest()[:16]
    digest = hashlib.sha1(f"{source_file_id}:{section_id}:{section_chunk_index}:{fingerprint}".encode("utf-8")).hexdigest()
    return f"chunk_{digest[:24]}"


def _link_neighbor_chunks(chunks: list[dict[str, Any]]) -> None:
    for index, chunk in enumerate(chunks):
        metadata = chunk.setdefault("metadata", {})
        if index > 0 and chunks[index - 1]["source_file_id"] == chunk["source_file_id"]:
            metadata["prev_chunk_id"] = chunks[index - 1]["id"]
        if index + 1 < len(chunks) and chunks[index + 1]["source_file_id"] == chunk["source_file_id"]:
            metadata["next_chunk_id"] = chunks[index + 1]["id"]


def _summary(text: str, max_chars: int = 320) -> str:
    normalized = " ".join(text.split())
    if len(normalized) <= max_chars:
        return normalized
    return normalized[:max_chars].rsplit(" ", 1)[0]


def _extract_equations(text: str) -> list[str]:
    equations = re.findall(r"\$\$(.+?)\$\$|\\\[(.+?)\\\]", text, flags=re.DOTALL)
    inline = re.findall(r"(?m)^\s*([A-Za-z0-9_(),+\-*/^=\s]{3,}=.+)$", text)
    extracted = [" ".join(part for part in match if part).strip() for match in equations]
    extracted.extend(item.strip() for item in inline if item.strip())
    return [item for item in extracted if item]


def _extract_tables(text: str) -> list[str]:
    tables: list[str] = []
    current: list[str] = []
    for line in text.splitlines():
        if line.strip().startswith("|") and line.strip().endswith("|"):
            current.append(line)
            continue
        if current:
            if len(current) >= 2:
                tables.append("\n".join(current))
            current = []
    if len(current) >= 2:
        tables.append("\n".join(current))
    return tables


def _extract_timeline_events(text: str) -> list[str]:
    event_re = re.compile(r"\b(?:\d{4}|Q[1-4]\s+\d{4}|week\s+\d+|module\s+\d+)\b", re.IGNORECASE)
    return [line.strip() for line in text.splitlines() if event_re.search(line)]


def _looks_like_chemical_equation(value: str) -> bool:
    return bool(
        re.search(r"(?:->|<-|<=>|=>|→|⇌|\\rightarrow|\\leftrightarrow)", value)
        and len(re.findall(r"(?:\d*[A-Z][a-z]?\d*)+", value)) >= 2
    )


def _extract_concepts(text: str) -> list[str]:
    candidates = re.findall(r"\b[A-Za-z][A-Za-z0-9_-]{4,}\b", text)
    stop = {"about", "these", "those", "which", "their", "there", "where", "student", "course"}
    counts: dict[str, int] = {}
    for candidate in candidates:
        key = candidate.lower()
        if key in stop:
            continue
        normalized = candidate[:1].upper() + candidate[1:]
        counts[normalized] = counts.get(normalized, 0) + 1
    return [item for item, _ in sorted(counts.items(), key=lambda pair: pair[1], reverse=True)[:8]]


def _generated_questions(concepts: list[str], heading_path: list[str], count: int) -> list[str]:
    topic = heading_path[-1] if heading_path else "this section"
    questions = [f"What does {topic} explain about {concept}?" for concept in concepts[:count]]
    while len(questions) < count:
        questions.append(f"What are the main ideas in {topic}?")
    return questions[:count]


_ingestion_service: IngestionService | None = None


def get_ingestion_service() -> IngestionService:
    global _ingestion_service
    if _ingestion_service is None:
        _ingestion_service = IngestionService()
    return _ingestion_service
